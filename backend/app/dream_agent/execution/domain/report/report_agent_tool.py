"""Report Agent Tool - BaseBizTool 기반 리포트 생성 도구

Markdown, HTML, PDF, PPTX 형식의 리포트를 생성합니다.
"""

import json
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, Optional
from string import Template

from backend.app.dream_agent.biz_execution.base_tool import (
    BaseBizTool,
    BizResult,
    BizResultStatus,
    BizResultMetadata,
    ApprovalType,
    ValidationResult
)
from backend.app.dream_agent.biz_execution.tool_registry import register_tool
from backend.app.dream_agent.models.todo import TodoItem


# ============================================================
# Report Templates
# ============================================================

MARKDOWN_TEMPLATE = """# ${title}

## 📊 분석 개요
- **분석일시**: ${timestamp}
- **브랜드**: ${brand}
- **채널**: ${channel}
- **총 리뷰 수**: ${total_reviews}
- **평균 평점**: ${average_rating} / 5.0

## 🎯 감성 분석 결과

### 전체 감성 분포
- 😊 **긍정**: ${positive_count}개 (${positive_ratio}%)
- 😐 **중립**: ${neutral_count}개
- 😞 **부정**: ${negative_count}개

## 💡 주요 인사이트

${insights}

## 📈 추천 사항

${recommendations}

## 📌 결론

${conclusion}

---
*본 보고서는 moaDREAM AI Agent에 의해 자동 생성되었습니다.*
"""

HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>${title}</title>
    <style>
        body { font-family: 'Pretendard', -apple-system, sans-serif; max-width: 900px; margin: 0 auto; padding: 20px; background: #f5f5f5; }
        .container { background: white; border-radius: 12px; padding: 40px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }
        h1 { color: #1a1a2e; border-bottom: 3px solid #e94560; padding-bottom: 10px; }
        h2 { color: #16213e; margin-top: 30px; }
        .stats { display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; margin: 20px 0; }
        .stat-card { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 20px; border-radius: 10px; text-align: center; }
        .stat-value { font-size: 2em; font-weight: bold; }
        .stat-label { opacity: 0.9; margin-top: 5px; }
        .sentiment { display: flex; gap: 15px; margin: 20px 0; }
        .sentiment-item { flex: 1; padding: 15px; border-radius: 8px; text-align: center; }
        .positive { background: #d4edda; color: #155724; }
        .neutral { background: #fff3cd; color: #856404; }
        .negative { background: #f8d7da; color: #721c24; }
        .insight-list { background: #f8f9fa; padding: 20px; border-radius: 8px; border-left: 4px solid #667eea; }
        .footer { text-align: center; color: #888; margin-top: 40px; padding-top: 20px; border-top: 1px solid #eee; }
    </style>
</head>
<body>
    <div class="container">
        <h1>📊 ${title}</h1>

        <div class="stats">
            <div class="stat-card">
                <div class="stat-value">${total_reviews}</div>
                <div class="stat-label">총 리뷰 수</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">${average_rating}</div>
                <div class="stat-label">평균 평점</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">${positive_ratio}%</div>
                <div class="stat-label">긍정 비율</div>
            </div>
        </div>

        <h2>🎯 감성 분석</h2>
        <div class="sentiment">
            <div class="sentiment-item positive">😊 긍정<br><strong>${positive_count}</strong></div>
            <div class="sentiment-item neutral">😐 중립<br><strong>${neutral_count}</strong></div>
            <div class="sentiment-item negative">😞 부정<br><strong>${negative_count}</strong></div>
        </div>

        <h2>💡 주요 인사이트</h2>
        <div class="insight-list">
            ${insights_html}
        </div>

        <h2>📈 추천 사항</h2>
        <div class="insight-list">
            ${recommendations_html}
        </div>

        <div class="footer">
            <p>moaDREAM AI Agent | ${timestamp}</p>
        </div>
    </div>
</body>
</html>
"""


# ============================================================
# Report Agent Tool
# ============================================================

@register_tool
class ReportAgentTool(BaseBizTool):
    """
    리포트 생성 도구

    분석 결과와 인사이트를 기반으로 다양한 형식의 리포트를 생성합니다.
    """

    name = "report_agent"
    description = "분석 결과 기반 리포트 생성 (Markdown, HTML, PDF, PPTX)"
    version = "2.0.0"

    requires_approval = False
    approval_type = ApprovalType.NONE

    is_async = False
    estimated_duration_sec = 30

    required_input_types = []  # 선택적 입력
    output_type = "report"

    has_cost = False

    def __init__(self):
        super().__init__()
        self.supported_formats = ["markdown", "html", "pdf", "pptx"]

    def validate_input(self, todo: TodoItem, context: Dict[str, Any]) -> ValidationResult:
        """입력 검증"""
        errors = []
        warnings = []

        # 포맷 검증
        format_param = todo.metadata.execution.tool_params.get("format", "markdown")
        if format_param not in self.supported_formats:
            errors.append(f"Unsupported format: {format_param}. Supported: {self.supported_formats}")

        # PDF/PPTX는 외부 라이브러리 필요 - 경고
        if format_param in ["pdf", "pptx"]:
            warnings.append(f"{format_param.upper()} format requires additional libraries. Will fallback to HTML if unavailable.")

        return ValidationResult(is_valid=len(errors) == 0, errors=errors, warnings=warnings)

    async def execute(
        self,
        todo: TodoItem,
        context: Dict[str, Any],
        log: Any
    ) -> BizResult:
        """리포트 생성 실행"""
        start_time = datetime.now()

        try:
            # 파라미터 추출
            params = todo.metadata.execution.tool_params
            format_type = params.get("format", "markdown")
            template_name = params.get("template", "default")

            # 분석 결과 로드
            analysis_data = self._load_analysis_data(context)
            insights_data = self._load_insights_data(context)

            # 리포트 데이터 준비
            report_data = self._prepare_report_data(analysis_data, insights_data, context)

            # 포맷별 생성
            if format_type == "markdown":
                content, output_path = self._generate_markdown(report_data)
            elif format_type == "html":
                content, output_path = self._generate_html(report_data)
            elif format_type == "pdf":
                content, output_path = self._generate_pdf(report_data)
            elif format_type == "pptx":
                content, output_path = self._generate_pptx(report_data)
            else:
                content, output_path = self._generate_markdown(report_data)

            # 처리 시간 계산
            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return self.create_result(
                todo=todo,
                status=BizResultStatus.SUCCESS,
                result_type="report",
                output_path=str(output_path),
                output_data={
                    "format": format_type,
                    "template": template_name,
                    "report_data": report_data
                },
                summary=f"{report_data['brand']} 분석 리포트 생성 완료 ({format_type.upper()})",
                preview=content[:500] + "..." if len(content) > 500 else content,
                metadata=BizResultMetadata(
                    processing_time_ms=processing_time,
                    output_size_bytes=len(content.encode('utf-8'))
                )
            )

        except Exception as e:
            return self.create_error_result(
                todo=todo,
                error_message=str(e),
                error_code="REPORT_GENERATION_ERROR"
            )

    def _load_analysis_data(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """분석 결과 로드"""
        # 컨텍스트에서 직접 로드
        if "analysis_result" in context:
            return context["analysis_result"]

        # 파일에서 로드 시도
        project_root = Path(__file__).parent.parent.parent.parent.parent.parent
        output_dir = project_root / "data/output/ml_results"

        if output_dir.exists():
            analysis_files = sorted(output_dir.glob("analysis_*.json"))
            if analysis_files:
                with open(analysis_files[-1], 'r', encoding='utf-8') as f:
                    return json.load(f)

        raise ValueError("analysis_result is required: provide it in context or ensure analysis file exists in data/output/ml_results/")

    def _load_insights_data(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """인사이트 로드"""
        if "insights" in context:
            return context["insights"]

        project_root = Path(__file__).parent.parent.parent.parent.parent.parent
        output_dir = project_root / "data/output/ml_results"

        if output_dir.exists():
            insight_files = sorted(output_dir.glob("insights_*.json"))
            if insight_files:
                with open(insight_files[-1], 'r', encoding='utf-8') as f:
                    return json.load(f)

        raise ValueError("insights data is required: provide it in context or ensure insights file exists in data/output/ml_results/")

    def _prepare_report_data(
        self,
        analysis: Dict[str, Any],
        insights: Dict[str, Any],
        context: Dict[str, Any]
    ) -> Dict[str, Any]:
        """리포트 데이터 준비"""
        sentiment = analysis.get("sentiment", {})

        return {
            "title": f"{context.get('brand', '브랜드')} 분석 리포트",
            "brand": context.get("brand", "K-Beauty 브랜드"),
            "channel": context.get("source", "올리브영"),
            "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "total_reviews": analysis.get("total_reviews", 0),
            "average_rating": analysis.get("average_rating", 0),
            "positive_count": sentiment.get("positive", 0),
            "neutral_count": sentiment.get("neutral", 0),
            "negative_count": sentiment.get("negative", 0),
            "positive_ratio": sentiment.get("positive_ratio", 0),
            "insights": insights.get("insights", []),
            "recommendations": insights.get("recommendations", []),
            "conclusion": insights.get("conclusion", "")
        }

    def _generate_markdown(self, data: Dict[str, Any]) -> tuple[str, Path]:
        """Markdown 리포트 생성"""
        template = Template(MARKDOWN_TEMPLATE)

        # 인사이트/추천 포맷팅
        insights_text = "\n".join(f"{i+1}. {insight}" for i, insight in enumerate(data["insights"]))
        recommendations_text = "\n".join(f"- {rec}" for rec in data["recommendations"])

        content = template.safe_substitute(
            title=data["title"],
            timestamp=data["timestamp"],
            brand=data["brand"],
            channel=data["channel"],
            total_reviews=data["total_reviews"],
            average_rating=data["average_rating"],
            positive_count=data["positive_count"],
            neutral_count=data["neutral_count"],
            negative_count=data["negative_count"],
            positive_ratio=data["positive_ratio"],
            insights=insights_text,
            recommendations=recommendations_text,
            conclusion=data["conclusion"]
        )

        # 파일 저장
        output_path = self._get_output_path("md")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(content, encoding="utf-8")

        return content, output_path

    def _generate_html(self, data: Dict[str, Any]) -> tuple[str, Path]:
        """HTML 리포트 생성"""
        template = Template(HTML_TEMPLATE)

        # HTML 리스트 포맷팅
        insights_html = "<ul>" + "".join(f"<li>{i}</li>" for i in data["insights"]) + "</ul>"
        recommendations_html = "<ul>" + "".join(f"<li>{r}</li>" for r in data["recommendations"]) + "</ul>"

        content = template.safe_substitute(
            title=data["title"],
            timestamp=data["timestamp"],
            total_reviews=data["total_reviews"],
            average_rating=data["average_rating"],
            positive_count=data["positive_count"],
            neutral_count=data["neutral_count"],
            negative_count=data["negative_count"],
            positive_ratio=data["positive_ratio"],
            insights_html=insights_html,
            recommendations_html=recommendations_html
        )

        output_path = self._get_output_path("html")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(content, encoding="utf-8")

        return content, output_path

    def _generate_pdf(self, data: Dict[str, Any]) -> tuple[str, Path]:
        """PDF 리포트 생성 (WeasyPrint 필요, 없으면 HTML 폴백)"""
        try:
            from weasyprint import HTML as WeasyprintHTML

            # HTML 먼저 생성
            html_content, _ = self._generate_html(data)

            output_path = self._get_output_path("pdf")
            output_path.parent.mkdir(parents=True, exist_ok=True)

            WeasyprintHTML(string=html_content).write_pdf(str(output_path))

            return f"PDF generated: {output_path}", output_path

        except ImportError:
            # WeasyPrint 없으면 HTML로 폴백
            return self._generate_html(data)

    def _generate_pptx(self, data: Dict[str, Any]) -> tuple[str, Path]:
        """PPTX 리포트 생성 (python-pptx 필요, 없으면 Markdown 폴백)"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # 타이틀 슬라이드
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = data["title"]
            subtitle.text = f"생성일: {data['timestamp']}"

            # 개요 슬라이드
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "분석 개요"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = f"총 리뷰: {data['total_reviews']}개"
            p = tf.add_paragraph()
            p.text = f"평균 평점: {data['average_rating']}"
            p = tf.add_paragraph()
            p.text = f"긍정 비율: {data['positive_ratio']}%"

            # 인사이트 슬라이드
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "주요 인사이트"
            body = slide.placeholders[1]
            tf = body.text_frame
            for i, insight in enumerate(data["insights"]):
                if i == 0:
                    tf.text = insight
                else:
                    p = tf.add_paragraph()
                    p.text = insight

            output_path = self._get_output_path("pptx")
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))

            return f"PPTX generated: {output_path}", output_path

        except ImportError:
            # python-pptx 없으면 Markdown으로 폴백
            return self._generate_markdown(data)

    def _get_output_path(self, extension: str) -> Path:
        """출력 경로 생성"""
        project_root = Path(__file__).parent.parent.parent.parent.parent.parent
        output_dir = project_root / "data/output/reports"
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return output_dir / f"report_{timestamp}.{extension}"
